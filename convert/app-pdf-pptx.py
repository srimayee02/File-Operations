import os
import sys
from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from flask import Flask, render_template, request, send_file

app = Flask(__name__)

def pdf_to_pptx(pdf_file):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # Create a working folder
    base_name = pdf_file.split(".pdf")[0]

    # Convert PDF to a list of images
    slideimgs = convert_from_path(pdf_file, 300, fmt='ppm', thread_count=2)

    # Loop over slides
    for i, slideimg in enumerate(slideimgs):
        if i % 10 == 0:
            print("Saving slide: " + str(i))

        imagefile = BytesIO()
        slideimg.save(imagefile, format='tiff')
        imagedata = imagefile.getvalue()
        imagefile.seek(0)
        width, height = slideimg.size

        # Set slide dimensions
        prs.slide_height = height * 9525
        prs.slide_width = width * 9525

        # Add slide
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(imagefile, 0, 0, width=width * 9525, height=height * 9525)

    # Save PowerPoint
    pptx_filename = base_name + '.pptx'
    prs.save(pptx_filename)
    return pptx_filename

@app.route('/', methods=['GET', 'POST'])
def convert_pdf_to_pptx():
    if request.method == 'POST' and 'pdf_file' in request.files:
        pdf_file = request.files['pdf_file']
        if pdf_file:
            pdf_filename = pdf_file.filename
            pdf_file.save(pdf_filename)

            # Convert PDF to PowerPoint
            pptx_filename = pdf_to_pptx(pdf_filename)

            # Provide the option to download the converted PowerPoint file
            return send_file(pptx_filename, as_attachment=True)

    return render_template('upload-pdf-pptx.html')

if __name__ == '__main__':
    app.run(debug=True)